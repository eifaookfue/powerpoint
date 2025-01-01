import sys
from typing import List
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph


def extract_text_from_pptx(file_path):
    # プレゼンテーションを読み込む
    presentation = Presentation(file_path)
    all_text: List[List[str]] = []

    # スライドごとにテキストを抽出
    for slide in presentation.slides:
        texts = [extract_text_from_shape(shape) for shape in slide.shapes]
        all_text.append(texts)

    return all_text


def extract_text_from_shape(shape: Shape) -> str:
    """Extract text from shape"""
    if not shape.has_text_frame:
        return None
    if shape.name == "PlaceHolder 1":
        return "\n".join([f"# {p.text}" for p in shape.text_frame.paragraphs])

    return "\n".join(
        [
            f"{create_bullet_point(p)}{p.text}"
            for p in shape.text_frame.paragraphs
            if p.text
        ]
    )


def create_bullet_point(paragraph: _Paragraph) -> str:
    """Create a bullet point according to paragraph.level"""
    return f"{'  ' * paragraph.level}- "


if __name__ == "__main__":
    args = sys.argv
    pptx_file = args[1]
    text_list = extract_text_from_pptx(pptx_file)

    for idx, texts in enumerate(text_list):
        for text in texts:
            print(text)
        if idx < len(text_list) - 1:
            print("---")
